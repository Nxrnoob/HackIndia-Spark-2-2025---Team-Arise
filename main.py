import os
import numpy as np
from sklearn.feature_extraction.text import TfidfVectorizer
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
import networkx as nx
import re

DOCUMENTS_DIR = os.path.join(os.path.dirname(__file__), "documents")

# Extract text from different file formats
def extract_text(file_name):
    file_path = os.path.join(DOCUMENTS_DIR, file_name)
    text = ""

    if not os.path.exists(file_path):
        return text

    if file_path.endswith(".pdf"):
        reader = PdfReader(file_path)
        text = "\n".join([page.extract_text() or "" for page in reader.pages])
    elif file_path.endswith(".docx"):
        doc = Document(file_path)
        text = "\n".join([para.text for para in doc.paragraphs])
    elif file_path.endswith(".pptx"):
        ppt = Presentation(file_path)
        text = "\n".join([shape.text for slide in ppt.slides for shape in slide.shapes if hasattr(shape, "text")])
    elif file_path.endswith(".txt"):
        with open(file_path, "r", encoding="utf-8") as f:
            text = f.read()

    return text.strip()

# Load all documents
def load_documents():
    docs = {}
    for file_name in os.listdir(DOCUMENTS_DIR):
        docs[file_name] = extract_text(file_name)
    return docs

# TextRank-based Summarization
def summarize_text(text, sentence_limit=3):
    sentences = re.split(r"(?<!\w\.\w.)(?<![A-Z][a-z]\.)(?<=\.|\?)\s", text)
    if len(sentences) <= sentence_limit:
        return text  # No need to summarize if it's short

    # Build similarity matrix
    vectorizer = TfidfVectorizer()
    sentence_vectors = vectorizer.fit_transform(sentences)
    similarity_matrix = np.dot(sentence_vectors, sentence_vectors.T).toarray()

    # Graph-based ranking (TextRank)
    nx_graph = nx.from_numpy_array(similarity_matrix)
    scores = nx.pagerank(nx_graph)
    ranked_sentences = sorted(((scores[i], s) for i, s in enumerate(sentences)), reverse=True)

    summary = " ".join([s[1] for s in ranked_sentences[:sentence_limit]])
    return summary

# TF-IDF Search + Context-Aware Recommendations
def search_documents(query, k=3):
    docs = load_documents()
    file_names = list(docs.keys())
    doc_texts = list(docs.values())

    if not doc_texts:
        return [{"file": "No relevant documents found.", "summary": ""}]

    vectorizer = TfidfVectorizer(stop_words='english')
    doc_vectors = vectorizer.fit_transform(doc_texts)
    query_vector = vectorizer.transform([query])

    similarities = np.dot(doc_vectors, query_vector.T).toarray().flatten()
    sorted_indices = np.argsort(similarities)[::-1]

    results = []
    relevant_keywords = set(query.lower().split())

    for idx in sorted_indices[:k]:
        if similarities[idx] > 0:
            file_name = file_names[idx]
            text = doc_texts[idx]
            summary = summarize_text(text)

            boost_score = 5 if "aiml" in file_name.lower() else 0
            total_score = similarities[idx] + boost_score

            results.append({
                "file": file_name,
                "summary": summary,
                "score": total_score
            })

    # Suggest Related Documents
    if results:
        main_doc_keywords = set(docs[results[0]["file"]].lower().split())
        related_docs = []
        for fname, text in docs.items():
            if fname != results[0]["file"]:
                overlap = len(main_doc_keywords.intersection(set(text.lower().split())))
                if overlap > 20:  # Adjust threshold based on content size
                    related_docs.append(fname)

        if related_docs:
            results.append({"file": "Related Documents:", "summary": ", ".join(related_docs)})

    return results if results else [{"file": "No relevant documents found.", "summary": ""}]

# Run search
if __name__ == "__main__":
    query = input("Enter search query: ")
    results = search_documents(query)

    print("\n🔍 **Search Results:**\n")
    for i, result in enumerate(results):
        print(f"{i+1}. **File:** {result['file']}\n   **Summary:** {result['summary']}\n")

